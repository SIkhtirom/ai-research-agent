"""Export service: build Markdown, PDF, and PPTX artefacts from a session's Q&A history."""

import io
from typing import Any

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
from sqlalchemy.orm import Session

from ...db.crud import QueryLogRepository, SessionRepository
from ...models import QueryLog, ResearchSession

SUPPORTED_EXPORT_FORMATS = ("md", "pdf", "pptx")


class ExportService:
    def export_session_data(
        self, db: Session, session_id: int, format_type: str
    ) -> dict[str, Any]:
        """Export a session's Q&A history into the requested format bytes."""
        if format_type not in SUPPORTED_EXPORT_FORMATS:
            raise ValueError(
                f"Unsupported format '{format_type}'. Allowed: {', '.join(SUPPORTED_EXPORT_FORMATS)}"
            )
        session = self.__load_session(db, session_id)
        query_logs = QueryLogRepository().list_by_session(db, session_id)

        if format_type == "md":
            content = self.__build_markdown(session, query_logs).encode("utf-8")
            return {"bytes": content, "filename": f"session_{session_id}.md", "media_type": "text/markdown"}

        if format_type == "pdf":
            content = self.__build_pdf(session, query_logs)
            return {"bytes": content, "filename": f"session_{session_id}.pdf", "media_type": "application/pdf"}

        content = self.__build_pptx(session, query_logs)
        return {"bytes": content, "filename": f"session_{session_id}.pptx", "media_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}

    def __load_session(self, db: Session, session_id: int) -> ResearchSession:
        session = SessionRepository().get_by_id(db, session_id)
        if session is None:
            raise ValueError(f"Session '{session_id}' does not exist.")
        return session

    def __build_markdown(
        self, session: ResearchSession, query_logs: list[QueryLog]
    ) -> str:
        sections = [f"# {session.title}", ""]
        if not query_logs:
            sections.append("_No conversations recorded in this session yet._")
        for index, query_log in enumerate(query_logs, start=1):
            sections.append(f"## Q{index}: {query_log.prompt}")
            sections.append("")
            sections.append(query_log.generated_response)
            sections.append("")
            sections.append("**Citations:**")
            for citation in query_log.citations_list:
                source_name = citation.get("source_name") or "unknown source"
                sections.append(f"- {source_name}")
            sections.append("")
        return "\n".join(sections)

    def __build_pdf(self, session: ResearchSession, query_logs: list[QueryLog]) -> bytes:
        buffer = io.BytesIO()
        document_builder = SimpleDocTemplate(
            buffer, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm
        )
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle("SessionTitle", parent=styles["Title"], fontSize=18)
        heading_style = ParagraphStyle("QuestionHead", parent=styles["Heading2"], fontSize=13)
        body_style = ParagraphStyle("AnswerBody", parent=styles["BodyText"], leading=14)
        citation_style = ParagraphStyle("CitationLine", parent=styles["Bullet"], fontSize=10)

        story = [Paragraph(self.__escape_markup(session.title), title_style), Spacer(1, 0.5 * cm)]
        if not query_logs:
            story.append(Paragraph("No conversations recorded in this session yet.", body_style))
        for index, query_log in enumerate(query_logs, start=1):
            story.append(Spacer(1, 0.4 * cm))
            story.append(Paragraph(f"Q{index}: {self.__escape_markup(query_log.prompt)}", heading_style))
            story.append(Spacer(1, 0.2 * cm))
            story.append(Paragraph(self.__escape_markup(query_log.generated_response), body_style))
            for citation in query_log.citations_list:
                source_name = citation.get("source_name") or "unknown source"
                story.append(Paragraph(f"&bull; {self.__escape_markup(str(source_name))}", citation_style))

        document_builder.build(story)
        return buffer.getvalue()

    def __build_pptx(self, session: ResearchSession, query_logs: list[QueryLog]) -> bytes:
        from pptx import Presentation

        presentation = Presentation()
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = session.title

        for query_log in query_logs:
            content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            content_slide.shapes.title.text = query_log.prompt
            body_placeholder = content_slide.placeholders[1]
            citation_lines = [
                f"- {citation.get('source_name') or 'unknown source'}"
                for citation in query_log.citations_list
            ]
            body_text = query_log.generated_response
            if citation_lines:
                body_text += "\n\nCitations:\n" + "\n".join(citation_lines)
            body_placeholder.text = body_text

        buffer = io.BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    def __escape_markup(self, raw_text: str) -> str:
        return (
            raw_text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )
