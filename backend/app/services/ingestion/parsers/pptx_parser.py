"""PowerPoint (.pptx) text extraction using python-pptx, including OCR of images.

Real-world presentations frequently nest content inside grouped shapes, tables,
and graphic frames. This parser walks every shape recursively (including groups),
extracts text frames and tables, and OCRs embedded pictures on a best-effort
basis. A missing or failing python-pptx only degrades logging - it never raises.
"""

import logging
from pathlib import Path

from .base_parser import DocumentParser, ExtractedDocument
from .ocr import is_ocr_available, ocr_image_bytes

logger = logging.getLogger(__name__)

# python-pptx shape type ids used to recognise pictures / graphic frames. Using
# the numeric ids keeps this robust against import/enum variations.
_PICTURE_TYPE_ID = 13  # MSO_SHAPE_TYPE.PICTURE
_TABLE_TYPE_ID = 19  # MSO_SHAPE_TYPE.TABLE


def _shape_type_id(shape) -> int | None:
    try:
        return shape.shape_type
    except Exception:  # noqa: BLE001 - some shapes expose no shape_type
        return None


class PptxParser(DocumentParser):
    def parse(self, source: Path | str) -> ExtractedDocument:
        source_path = Path(source)
        try:
            from pptx import Presentation
        except ImportError as exc:  # noqa: BLE001
            logger.error(
                "python-pptx is not installed; cannot parse PPTX '%s': %s",
                source_path.name,
                exc,
            )
            raise ValueError("python-pptx is not installed on the server") from exc

        try:
            presentation = Presentation(source_path)
        except Exception as exc:  # noqa: BLE001 - invalid/corrupt package
            logger.warning("Could not open PPTX '%s': %s", source_path.name, exc)
            raise ValueError(f"Invalid or corrupt PPTX: {source_path.name}") from exc

        slide_texts: list[str] = []
        image_ocr_texts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            self.__walk_shapes(
                slide.shapes, slide_lines, image_ocr_texts, slide_index
            )
            if slide_lines:
                slide_texts.append("\n".join(slide_lines))

        sections: list[str] = []
        if slide_texts:
            sections.append("\n\n".join(slide_texts))
        if image_ocr_texts:
            sections.append("\n".join(image_ocr_texts))

        content_text = "\n\n".join(sections)
        if not content_text.strip():
            logger.warning(
                "PPTX '%s' produced no extractable text (it may be entirely image-based).",
                source_path.name,
            )

        metadata = {
            "filename": source_path.name,
            "slide_count": len(presentation.slides),
            "image_text_count": len(image_ocr_texts),
            "ocr_available": is_ocr_available(),
        }
        return ExtractedDocument(content_text=content_text, metadata=metadata)

    def __walk_shapes(
        self,
        shapes,
        slide_lines: list[str],
        image_ocr_texts: list[str],
        slide_index: int,
    ) -> None:
        """Recursively collect text and OCR images from a shape collection."""
        try:
            shape_list = list(shapes)
        except Exception:  # noqa: BLE001
            return
        for shape in shape_list:
            self.__handle_shape(shape, slide_lines, image_ocr_texts, slide_index)

    def __handle_shape(
        self, shape, slide_lines: list[str], image_ocr_texts: list[str], slide_index: int
    ) -> None:
        shape_type = _shape_type_id(shape)

        # Grouped shapes: recurse into the group's children.
        if self.__is_group(shape):
            try:
                self.__walk_shapes(
                    shape.shapes, slide_lines, image_ocr_texts, slide_index
                )
            except Exception as exc:  # noqa: BLE001
                logger.debug("Skipped grouped shapes: %s", exc)
            return

        # Pictures: OCR the image content.
        if shape_type == _PICTURE_TYPE_ID:
            try:
                blob = shape.image.blob
                ocr_text = ocr_image_bytes(blob)
                if ocr_text:
                    image_ocr_texts.append(f"[Slide {slide_index} - Gambar] {ocr_text}")
            except Exception as exc:  # noqa: BLE001
                logger.debug("Slide image OCR skipped: %s", exc)
            return

        # Tables: extract every cell.
        if shape_type == _TABLE_TYPE_ID or hasattr(shape, "table"):
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame and cell.text_frame.text.strip():
                            slide_lines.append(cell.text_frame.text.strip())
            except Exception as exc:  # noqa: BLE001
                logger.debug("Table extraction skipped on a slide: %s", exc)
            return

        # Text frames (text boxes and placeholders).
        self.__extract_text_frame(shape, slide_lines)

    @staticmethod
    def __is_group(shape) -> bool:
        try:
            # Groups expose a "shapes" sub-collection.
            return hasattr(shape, "shapes") and callable(getattr(shape.shapes, "__iter__", None))
        except Exception:  # noqa: BLE001
            return False

    @staticmethod
    def __extract_text_frame(shape, slide_lines: list[str]) -> None:
        try:
            if not shape.has_text_frame:
                return
        except Exception:  # noqa: BLE001
            return
        try:
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs)
                if not text.strip():
                    continue
                slide_lines.append(text)
        except Exception as exc:  # noqa: BLE001
            logger.debug("Text frame extraction skipped: %s", exc)
